"""
Extracts text content from PPT (PPTX) and PDF files for AI analysis.
"""
import io
import os
from typing import Optional


def parse_pptx(file_bytes: bytes) -> dict:
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed", "slides": [], "full_text": ""}
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        all_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            title = ""
            slide_title = getattr(slide.shapes, "title", None)

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)
                if slide_title is not None and shape == slide_title and hasattr(shape, "text"):
                    title = shape.text.strip()

            content = "\n".join(slide_text)
            slides.append({
                "slide_number": i,
                "title": title or f"Slide {i}",
                "content": content,
            })
            all_text.append(content)

        return {
            "slides": slides,
            "full_text": "\n\n".join(all_text),
            "total_slides": len(slides),
            "format": "pptx",
        }
    except Exception as exc:
        return {"error": f"Could not read PPTX file: {exc}", "slides": [], "full_text": ""}


def parse_pdf(file_bytes: bytes) -> dict:
    """Extract text from a PDF file."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return {"error": "PyPDF2 not installed", "slides": [], "full_text": ""}
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        slides = []
        all_text = []

        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            text = text.strip()
            lines = text.split("\n")
            title = lines[0] if lines else f"Page {i}"

            slides.append({
                "slide_number": i,
                "title": title[:100],
                "content": text,
            })
            all_text.append(text)

        return {
            "slides": slides,
            "full_text": "\n\n".join(all_text),
            "total_slides": len(slides),
            "format": "pdf",
        }
    except Exception as exc:
        return {"error": f"Could not read PDF file: {exc}", "slides": [], "full_text": ""}


def parse_file(file_bytes: bytes, filename: str) -> dict:
    """Auto-detect format and parse."""
    if not filename:
        return {"error": "Missing filename", "slides": [], "full_text": ""}

    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pptx":
        return parse_pptx(file_bytes)
    elif ext == ".pdf":
        return parse_pdf(file_bytes)
    else:
        return {"error": f"Unsupported format: {ext}", "slides": [], "full_text": ""}
